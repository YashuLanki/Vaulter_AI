"""
pipeline/email_reader.py
-------------------------
Vaulter AI Stage 2 — Outlook Email Pipeline

Reads broker and market-report emails from Outlook via Microsoft Graph,
extracts clean text, and stores it in the same ChromaDB as Stage 1 PDFs.

PDF attachments are routed to data/watched_folder/ so Stage 1 picks
them up automatically.

Called by:  python main.py email
            python main.py email --days 30
            pipeline/scheduler.py (on a timer)
"""

import hashlib
import json
import logging
import re
import shutil
import sys
from datetime import datetime, timedelta, timezone
from pathlib import Path

import requests

sys.path.insert(0, str(Path(__file__).parent.parent))
import safe_io
from config import (
    LOG_DIR, RAW_EMAIL_DIR, DATA_DIR, REGISTRY_DIR,
    WATCH_DIR, CHROMA_DIR, CHROMA_COLLECTION_NAME,
    LOG_LEVEL,
    OUTLOOK_FOLDERS, OUTLOOK_SENDER_WHITELIST, OUTLOOK_LOOKBACK_DAYS,
)
from pipeline.outlook_auth import get_access_token
from pipeline.property_matcher import match_properties, matched_property_tags, format_matched_properties

# ─── Logging ──────────────────────────────────────────────────────
LOG_DIR.mkdir(parents=True, exist_ok=True)
logging.basicConfig(
    level=getattr(logging, LOG_LEVEL, logging.INFO),
    format="%(asctime)s [EMAIL] %(levelname)s — %(message)s",
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler(LOG_DIR / "email_reader.log", encoding="utf-8"),
    ],
)
log = logging.getLogger(__name__)

GRAPH = "https://graph.microsoft.com/v1.0"

# ─── Registry (tracks processed message IDs) ──────────────────────
REGISTRY_FILE = REGISTRY_DIR / "email_registry.json"

def load_registry() -> set:
    return set(safe_io.load_json(REGISTRY_FILE, default=[]))

def save_registry(seen: set):
    """Merges `seen` into whatever's currently on disk (under a file
    lock) via set union, rather than overwriting -- process_all_emails()
    accumulates seen_ids in memory for the whole run and saves once at
    the end, so a blind overwrite here would discard any message ID
    another concurrently-running email check (e.g. a manual run
    overlapping the scheduler) already recorded as seen."""
    safe_io.locked_json_update(
        REGISTRY_FILE, lambda current: sorted(set(current) | seen), default=[],
    )


# ─── Graph API ────────────────────────────────────────────────────

def graph_get(token: str, path: str, params: dict = None) -> dict | None:
    headers = {"Authorization": f"Bearer {token}", "Accept": "application/json"}
    try:
        resp = requests.get(
            f"{GRAPH}{path}", headers=headers, params=params, timeout=20
        )
        resp.raise_for_status()
        return resp.json()
    except requests.HTTPError as e:
        log.error(f"Graph error [{resp.status_code}] {path}: {e}")
    except Exception as e:
        log.error(f"Request failed {path}: {e}")
    return None

def get_folder_id(token: str, folder_name: str) -> str | None:
    """Resolve folder display name to Graph folder ID."""
    well_known = {
        "inbox": "inbox", "drafts": "drafts", "sentitems": "sentItems",
        "deleteditems": "deletedItems", "junkemail": "junkemail",
    }
    if folder_name.lower() in well_known:
        return well_known[folder_name.lower()]

    data = graph_get(token, "/me/mailFolders", {"$filter": f"displayName eq '{folder_name}'"})
    if data and data.get("value"):
        return data["value"][0]["id"]

    log.warning(f"Outlook folder '{folder_name}' not found — skipping")
    return None

def list_messages(token: str, folder_id: str, lookback_days: int) -> list[dict]:
    """Only lists messages received within the last lookback_days — without
    this $filter, Graph returns the entire folder every call, ignoring the
    parameter entirely and re-paging the whole mailbox on every scheduled run."""
    cutoff = (datetime.now(timezone.utc) - timedelta(days=lookback_days)).strftime("%Y-%m-%dT%H:%M:%SZ")
    params = {
        "$select": "id,subject,from,receivedDateTime",
        "$filter": f"receivedDateTime ge {cutoff}",
        "$top": 50,
        "$orderby": "receivedDateTime desc",
    }
    messages, path = [], f"/me/mailFolders/{folder_id}/messages"
    while path:
        data = graph_get(token, path.replace(GRAPH, ""), params)
        if not data:
            break
        messages.extend(data.get("value", []))
        path   = data.get("@odata.nextLink")
        params = None
    return messages

def get_body(token: str, msg_id: str) -> str | None:
    """Returns the cleaned body text, or None if the Graph fetch itself
    failed (expired token, network error, etc.) — distinct from a
    genuinely short/empty body, which is a legitimate "" result. Callers
    must NOT treat a None return the same as a real empty body, or a
    transient fetch failure gets mistaken for "nothing to store" and the
    email is lost forever once marked seen."""
    data = graph_get(token, f"/me/messages/{msg_id}", {"$select": "body"})
    if data is None:
        return None
    body    = data.get("body", {})
    content = body.get("content", "")
    if body.get("contentType", "").lower() == "html":
        # Strip HTML tags
        content = re.sub(r"<[^>]+>", " ", content)
        content = re.sub(r"&nbsp;", " ", content)
        content = re.sub(r"&amp;", "&", content)
        content = re.sub(r"&lt;", "<", content)
        content = re.sub(r"&gt;", ">", content)
        content = re.sub(r"&quot;", '"', content)
    return clean_email_body(" ".join(content.split()))


def clean_email_body(text: str) -> str:
    """
    Strip signatures, legal disclaimers, and forwarded/replied-to headers
    from an email body so only the actual message content is stored.

    Removes:
      - Outlook/Gmail signature blocks (-- / ___ dividers)
      - Legal disclaimers (confidentiality notices, CAUTION headers)
      - Forwarded message headers (From: ... Sent: ... To: ... Subject:)
      - Replied-to quote blocks (lines starting with >)
      - Excessive whitespace
    """
    lines = text.splitlines()
    cleaned = []

    # Patterns that signal the start of a signature or disclaimer block
    STOP_PATTERNS = [
        r"^--\s*$",                                    # -- signature divider
        r"^_{3,}\s*$",                                 # ___ divider
        r"^-{3,}\s*$",                                 # --- divider
        r"^\*{3,}\s*$",                               # *** divider
        r"(?i)^(CAUTION|DISCLAIMER|CONFIDENTIAL)",      # legal headers
        r"(?i)this (e-?mail|message) (and any|is (confidential|intended))",
        r"(?i)if you (have received|are not the intended)",
        r"(?i)^(From|Sent|To|Cc|Subject):\s+",         # forwarded header block
        r"(?i)-----\s*original message\s*-----",
        r"(?i)-----\s*forwarded message\s*-----",
        r"(?i)^on .{10,} wrote:$",                      # "On Mon, Jan 1 John wrote:"
        r"(?i)get outlook for",                          # Outlook mobile footer
        r"(?i)sent from my (iphone|ipad|android|samsung|mobile)",
    ]

    import re as _re
    stop_patterns = [_re.compile(p) for p in STOP_PATTERNS]

    for line in lines:
        stripped = line.strip()

        # Skip quoted reply lines
        if stripped.startswith(">"):
            continue

        # Check if this line starts a signature/disclaimer block — stop here
        if any(pat.search(stripped) for pat in stop_patterns):
            break

        cleaned.append(line)

    result = "\n".join(cleaned).strip()

    # Collapse runs of blank lines to a single blank line
    result = _re.sub(r"\n{3,}", "\n\n", result)

    # If cleaning removed too much, fall back to the raw text
    if len(result) < 30 and len(text) > 30:
        return text

    return result

def get_attachments(token: str, msg_id: str) -> list[dict] | None:
    """Returns the attachment list, or None if the Graph fetch itself
    failed (expired token, network error, etc.) — distinct from [], which
    means the message genuinely has no attachments. Callers must NOT
    treat a None return the same as a real empty list, or a transient
    fetch failure gets mistaken for "no attachments" and whatever was
    actually attached (a survey, a DD report) is never fetched again."""
    data = graph_get(token, f"/me/messages/{msg_id}/attachments")
    if data is None:
        return None
    return [
        {"id": a["id"], "name": a.get("name", "attachment"),
         "mime_type": a.get("contentType", ""), "size": a.get("size", 0)}
        for a in data.get("value", [])
        if not a.get("isInline", False)
    ]

def get_attachment_bytes(token: str, msg_id: str, att_id: str) -> bytes:
    try:
        resp = requests.get(
            f"{GRAPH}/me/messages/{msg_id}/attachments/{att_id}/$value",
            headers={"Authorization": f"Bearer {token}"}, timeout=30,
        )
        resp.raise_for_status()
        return resp.content
    except Exception as e:
        log.error(f"Attachment download failed: {e}")
        return b""


# ─── ChromaDB ─────────────────────────────────────────────────────

def get_collection():
    from ingestion.embedder import get_collection as _get_collection
    return _get_collection()

def simple_embed(text: str) -> list[float]:
    from ingestion.embedder import embed_texts
    return embed_texts([text])[0]

def chunk_text(text: str, page_count: int = 1) -> list[str]:
    from ingestion.chunker import chunk_text as _chunk_text
    return _chunk_text(text, page_count=page_count)

def store_email(msg_id, subject, sender, date_str, body, collection):
    if len(body) < 50:
        log.info(f"Body too short ({len(body)} chars) — skipping")
        return
    # Match the email body text to relevant properties
    from pipeline.property_matcher import match_properties, matched_property_tags, format_matched_properties
    search_text = f"{subject} {body}"
    matches     = match_properties(search_text)
    prop_tags   = matched_property_tags(matches)
    if matches:
        log.info(f"  Email matched: {format_matched_properties(matches)}")

    # Build a short preview of the message for Claude to use in summaries
    # Takes the first meaningful sentence(s) up to 300 chars
    body_preview = body[:300].rsplit(" ", 1)[0] + "..." if len(body) > 300 else body

    chunks = chunk_text(body)
    ids, docs, metas, embeds = [], [], [], []
    for i, chunk in enumerate(chunks):
        ids.append(f"email_{msg_id}_{i}")
        docs.append(chunk)
        metas.append({
            "source":       sender,
            "subject":      subject[:200],
            "date":         date_str,
            "msg_id":       msg_id,
            "chunk":        i,
            "type":         "email",
            "body_preview": body_preview,
            **prop_tags,
        })
        embeds.append(simple_embed(chunk))
    collection.upsert(ids=ids, documents=docs, metadatas=metas, embeddings=embeds)
    log.info(f"  ✓ {len(chunks)} chunks stored | '{subject[:60]}'")


# ─── Attachment Routing ───────────────────────────────────────────

def _store_attachment_chunks(name, msg_id, subject, text, file_type, collection,
                             prop_tags=None, page_count: int = 1):
    """Helper — chunks text and stores in ChromaDB with attachment metadata."""
    chunks = chunk_text(text, page_count=page_count)
    if not chunks:
        return 0
    ids, docs, metas, embeds = [], [], [], []
    for i, chunk in enumerate(chunks):
        cid = f"email_att_{hashlib.sha256(chunk.encode()).hexdigest()[:12]}_{i}"
        ids.append(cid)
        docs.append(chunk)
        metas.append({
            "source":   name,
            "filename": name,
            "subject":  subject[:200],
            "msg_id":   msg_id,
            "chunk":    i,
            "type":     f"email_attachment_{file_type}",
            **(prop_tags or {}),
        })
        embeds.append(simple_embed(chunk))
    collection.upsert(ids=ids, documents=docs, metadatas=metas, embeddings=embeds)
    return len(chunks)


def _match_and_log(name, subject, text) -> tuple:
    """Run property matcher and log result. Returns (prop_tags, matches)."""
    from pipeline.property_matcher import match_from_filename

    # File routing is based on filename only — subject/text match is for
    # ChromaDB tagging only, not for deciding which folder to save to.
    # A generic filename like CostarExport.xlsx should go to general/
    # even if the subject mentions a city.
    file_matches = match_from_filename(name)
    strong_file_matches = [m for m in file_matches
                           if "name_match" in m.get("match_reasons", [])
                           or "city_match" in m.get("match_reasons", [])]

    # For ChromaDB tags — use full subject+text match
    all_matches  = match_properties(f"{subject} {text[:500]}")
    prop_tags    = matched_property_tags(all_matches)
    matched_desc = format_matched_properties(all_matches)
    log.info(f"  Matched properties: {matched_desc}")

    # Return strong file matches for routing, all matches for ChromaDB tags
    return prop_tags, strong_file_matches


def _save_to_processed(raw_path: Path, name: str, prop_tags: dict, ts: str,
                       matches: list = None):
    """
    Copy an attachment to processed/State/Property/ so it is accessible
    via File Explorer.

    Only routes to a property folder on strong matches (name_match or city_match).
    Weak matches (state_match, category_match only) go to processed/general/
    so files are not placed in the wrong property folder.
    """
    try:
        from config import PROCESSED_DIR

        # Check if we have a strong match. IMPORTANT: derive state/prop from
        # THIS SAME strong match, not from prop_tags -- prop_tags comes from
        # a completely different (subject+body text) match pass, and the two
        # can disagree (e.g. a filename that clearly names a property, on a
        # generic-subject broker email that never mentions it in the body).
        # Using prop_tags here previously meant that disagreement could
        # produce empty state/prop strings, which silently collapsed the
        # destination to PROCESSED_DIR's own root -- not the property
        # folder, and not processed/general/ either.
        strong_match  = None
        if matches:
            for m in matches:
                reasons = m.get("match_reasons", [])
                if "name_match" in reasons or "city_match" in reasons:
                    strong_match = m
                    break

        if strong_match:
            state    = strong_match.get("state", "").strip().lower().replace(" ", "_")
            prop     = strong_match.get("name", "").strip()
            dest_dir = PROCESSED_DIR / state / prop
            log.info(f"  Saved to processed: {state}/{prop}/{name}")
        else:
            dest_dir = PROCESSED_DIR / "general"
            log.info(f"  Saved to processed/general/ (no strong property match): {name}")

        dest_dir.mkdir(parents=True, exist_ok=True)
        # shutil.copy silently overwrites a same-named file already sitting
        # here (e.g. two different emails both attaching "invoice.pdf" for
        # the same property) -- dedupe_path avoids destroying the earlier one.
        dest = safe_io.dedupe_path(dest_dir, name)
        shutil.copy(str(raw_path), str(dest))

    except Exception as e:
        log.debug(f"  Could not save to processed: {e}")


def route_attachment(token, msg_id, att, subject, collection):
    """
    Route email attachments by file type:
      PDF          → Stage 1 watched_folder (OCR handled there)
      Word (.docx) → mammoth text extraction
      Excel (.xlsx)→ openpyxl text extraction
      PowerPoint   → python-pptx text extraction
      Images       → Tesseract OCR
      CSV/Text     → direct text read
      Other        → saved to raw_email for manual review
    """
    name      = att["name"]
    mime_type = att["mime_type"].lower()
    data      = get_attachment_bytes(token, msg_id, att["id"])
    if not data:
        # Download failure is indistinguishable from a genuinely empty
        # attachment here (get_attachment_bytes returns b"" for both). This
        # attachment silently won't be stored -- log it so a run of failed
        # downloads is at least visible, rather than a silent no-op.
        log.warning(f"  Attachment download failed or empty, skipping: {name}")
        return

    RAW_EMAIL_DIR.mkdir(parents=True, exist_ok=True)
    ts        = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_name = re.sub(r"[^\w\-.]", "_", name)
    raw_path  = RAW_EMAIL_DIR / f"{ts}_{safe_name}"
    raw_path.write_bytes(data)
    nl        = name.lower()

    # ── PDF → Stage 1 watched_folder (routed to State/Property subfolder) ──
    if "pdf" in mime_type or nl.endswith(".pdf"):
        # Try to match the PDF to a property so it lands in the right subfolder
        try:
            from pipeline.property_matcher import match_from_filename, match_properties
            # 1. Try filename first
            matches = match_from_filename(name)
            # 2. Fall back to subject + filename
            if not matches:
                matches = match_properties(f"{subject} {name}")
            # 3. Last resort: extract text from the PDF itself and match against content
            if not matches:
                try:
                    import pdfplumber
                    with pdfplumber.open(str(raw_path)) as pdf:
                        # Read first 5 pages — enough to identify the property
                        pages_to_read = min(5, len(pdf.pages))
                        pdf_text = ""
                        for page in pdf.pages[:pages_to_read]:
                            t = page.extract_text()
                            if t:
                                pdf_text += t + " "
                    if pdf_text.strip():
                        matches = match_properties(pdf_text[:5000])  # first 5000 chars
                        if matches:
                            log.info(f"  Matched via PDF content: {matches[0]['name']}")
                except Exception as pdf_err:
                    log.debug(f"  PDF content read failed: {pdf_err}")
            if matches:
                best = matches[0]
                prop_state = best.get("state", "").strip()
                prop_name  = best.get("name", "").strip()
                dest_dir   = WATCH_DIR / prop_state / prop_name
            else:
                # No property match — route to processed/general/ so the
                # open_general_files MCP tool can surface it in claude.ai
                from config import PROCESSED_DIR
                dest_dir = PROCESSED_DIR / "general"
                log.info(f"  PDF attachment unmatched — routed to processed/general/")
        except Exception:
            from config import PROCESSED_DIR
            dest_dir = PROCESSED_DIR / "general"
            log.info(f"  PDF attachment routing error — routed to processed/general/")

        dest_dir.mkdir(parents=True, exist_ok=True)
        # Keyed by content hash, not the fetch timestamp -- if this same
        # email/attachment is ever processed more than once (a registry
        # reset, an overlapping lookback window, a re-run after a crash),
        # a timestamp-keyed name would land as a brand-new, uniquely-named
        # file every time. Stage 1's own hash-based dedup would correctly
        # skip re-ingesting it, but skipped files are never moved out of
        # watched_folder -- so the folder would silently accumulate one
        # never-cleaned-up duplicate copy per reprocessing. A stable,
        # content-derived name means reprocessing overwrites the exact
        # same file instead of piling up new ones.
        content_key = hashlib.sha256(data).hexdigest()[:12]
        dest_file   = dest_dir / f"email_{content_key}_{safe_name}"
        shutil.copy(raw_path, dest_file)
        log.info(f"  PDF → {dest_dir.relative_to(WATCH_DIR)}/: {safe_name}")

    # ── Word Document ──────────────────────────────────────────────
    elif nl.endswith((".docx", ".doc")) or "wordprocessingml" in mime_type:
        try:
            import mammoth
            with open(raw_path, "rb") as f:
                text = " ".join(mammoth.extract_raw_text(f).value.split())
            if len(text) < 50:
                log.info(f"  Word doc too short: {name}")
                return
            prop_tags, matches = _match_and_log(name, subject, text)
            n = _store_attachment_chunks(name, msg_id, subject, text, "docx",
                                         collection, prop_tags)
            _save_to_processed(raw_path, name, prop_tags, ts, matches)
            log.info(f"  ✓ Word doc stored: {name} ({n} chunks)")
        except ImportError:
            log.warning(f"  mammoth not installed — pip install mammoth")
        except Exception as e:
            log.error(f"  Word doc failed {name}: {e}")

    # ── Excel Spreadsheet ──────────────────────────────────────────
    elif nl.endswith((".xlsx", ".xlsm", ".xls")) or "spreadsheetml" in mime_type:
        try:
            import openpyxl, io
            wb    = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
            lines = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        lines.append(row_text)
            # Double newline between rows so the chunker treats each row as its
            # own paragraph and keeps it intact. Single \n made the entire sheet
            # one giant paragraph, causing the chunker to cut rows at arbitrary
            # 500-char boundaries — fragmenting pipe-separated data.
            # page_count=9999 triggers the 8000-char chunk tier in config.CHUNK_TIERS,
            # which is large enough to hold any Excel row without splitting it.
            # This only affects Excel files — all other attachment types use the
            # default page_count=1 (500-char chunks) and are unaffected.
            text = "\n\n".join(lines)
            if len(text) < 50:
                log.info(f"  Excel too short: {name}")
                return
            prop_tags, matches = _match_and_log(name, subject, text)
            n = _store_attachment_chunks(name, msg_id, subject, text, "excel",
                                         collection, prop_tags, page_count=9999)
            _save_to_processed(raw_path, name, prop_tags, ts, matches)
            log.info(f"  ✓ Excel stored: {name} ({n} chunks)")
        except ImportError:
            log.warning(f"  openpyxl not installed — pip install openpyxl")
        except Exception as e:
            log.error(f"  Excel failed {name}: {e}")

    # ── PowerPoint ────────────────────────────────────────────────
    elif nl.endswith((".pptx", ".ppt")) or "presentationml" in mime_type:
        try:
            from pptx import Presentation
            import io
            prs   = Presentation(io.BytesIO(data))
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text.strip())
            text = " ".join(lines)
            if len(text) < 50:
                log.info(f"  PowerPoint too short: {name}")
                return
            prop_tags, matches = _match_and_log(name, subject, text)
            n = _store_attachment_chunks(name, msg_id, subject, text, "pptx",
                                         collection, prop_tags)
            _save_to_processed(raw_path, name, prop_tags, ts, matches)
            log.info(f"  ✓ PowerPoint stored: {name} ({n} chunks)")
        except ImportError:
            log.warning(f"  python-pptx not installed — pip install python-pptx")
        except Exception as e:
            log.error(f"  PowerPoint failed {name}: {e}")

    # ── CSV ───────────────────────────────────────────────────────
    elif nl.endswith(".csv") or "text/csv" in mime_type:
        try:
            import csv, io
            text_data = data.decode("utf-8", errors="replace")
            reader    = csv.reader(io.StringIO(text_data))
            lines     = [" | ".join(row) for row in reader if any(c.strip() for c in row)]
            text      = "\n".join(lines)
            if len(text) < 50:
                log.info(f"  CSV too short: {name}")
                return
            prop_tags, matches = _match_and_log(name, subject, text)
            n = _store_attachment_chunks(name, msg_id, subject, text, "csv",
                                         collection, prop_tags)
            log.info(f"  ✓ CSV stored: {name} ({n} chunks)")
        except Exception as e:
            log.error(f"  CSV failed {name}: {e}")

    # ── Images → OCR ─────────────────────────────────────────────
    elif nl.endswith((".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp")):
        try:
            import pytesseract
            from PIL import Image
            import io
            from config import TESSERACT_PATH
            if TESSERACT_PATH:
                pytesseract.pytesseract.tesseract_cmd = str(TESSERACT_PATH)
            img  = Image.open(io.BytesIO(data))
            text = " ".join(pytesseract.image_to_string(img).split())
            if len(text) < 50:
                log.info(f"  Image OCR too short: {name}")
                return
            prop_tags, matches = _match_and_log(name, subject, text)
            n = _store_attachment_chunks(name, msg_id, subject, text, "image_ocr",
                                         collection, prop_tags)
            log.info(f"  ✓ Image OCR stored: {name} ({n} chunks)")
        except ImportError:
            log.warning(f"  pytesseract/PIL not installed")
        except Exception as e:
            log.error(f"  Image OCR failed {name}: {e}")

    # ── Plain Text ────────────────────────────────────────────────
    elif "text" in mime_type or nl.endswith(".txt"):
        text = data.decode("utf-8", errors="replace")
        prop_tags, matches = _match_and_log(name, subject, text)
        n = _store_attachment_chunks(name, msg_id, subject, text, "text",
                                     collection, prop_tags)
        log.info(f"  ✓ Text stored: {name} ({n} chunks)")

    # ── Unsupported ───────────────────────────────────────────────
    else:
        log.info(f"  Attachment skipped (unsupported type): {name} [{att['mime_type']}]")



# ─── Main ─────────────────────────────────────────────────────────

def _sender_is_whitelisted(sender: str, whitelist: list[str]) -> bool:
    """
    Match a sender address against OUTLOOK_SENDER_WHITELIST.

    Deliberately NOT a substring check (`entry in sender`) -- that would
    let a lookalike/spoofed domain like "broker@trusted-partner.com.evil.net"
    match a whitelist entry of "trusted-partner.com", since the real
    domain is just a substring of the fake one. A whitelist entry that
    contains "@" must match the sender's full address exactly; a bare
    domain entry must equal the sender's domain exactly (not merely be
    contained in it).
    """
    sender_l = sender.strip().lower()
    if "@" not in sender_l:
        return False
    sender_domain = sender_l.rsplit("@", 1)[1]

    for entry in whitelist:
        entry_l = entry.strip().lower()
        if "@" in entry_l:
            if sender_l == entry_l:
                return True
        elif sender_domain == entry_l:
            return True
    return False


def process_all_emails(lookback_days: int | None = None):
    RAW_EMAIL_DIR.mkdir(parents=True, exist_ok=True)

    if lookback_days is None:
        lookback_days = OUTLOOK_LOOKBACK_DAYS

    token      = get_access_token()
    collection = get_collection()
    seen_ids   = load_registry()

    new = skipped = errors = 0

    for folder_name in OUTLOOK_FOLDERS:
        folder_id = get_folder_id(token, folder_name)
        if not folder_id:
            continue

        messages = list_messages(token, folder_id, lookback_days)
        log.info(f"Folder '{folder_name}': {len(messages)} messages in window")

        for msg in messages:
            msg_id   = msg["id"]
            subject  = msg.get("subject") or "(no subject)"
            sender   = msg.get("from", {}).get("emailAddress", {}).get("address", "unknown")
            date_str = msg.get("receivedDateTime", "")

            if msg_id in seen_ids:
                skipped += 1
                continue

            if OUTLOOK_SENDER_WHITELIST:
                if not _sender_is_whitelisted(sender, OUTLOOK_SENDER_WHITELIST):
                    # Deliberately NOT marked seen -- this message was never
                    # actually processed, just rejected by the current
                    # whitelist. If the whitelist is later updated to
                    # include this sender (e.g. a new broker is onboarded),
                    # this email needs to still be eligible for processing
                    # on the next run instead of being permanently lost.
                    # list_messages() already fetched this message's
                    # metadata in one bulk call for the whole folder, so
                    # re-checking it against the whitelist every run costs
                    # nothing extra -- only actually-processed messages
                    # need the "don't refetch" guarantee seen_ids provides.
                    skipped += 1
                    continue

            log.info(f"Processing: '{subject[:70]}' from {sender}")
            try:
                body = get_body(token, msg_id)
                if body is None:
                    # Fetch failed (expired token, network error, etc.) --
                    # do NOT mark seen, so this message is retried next run
                    # instead of being permanently lost with no content stored.
                    log.error(f"  Could not fetch body for '{subject[:50]}' — will retry next run")
                    errors += 1
                    continue

                # Audit trail
                (RAW_EMAIL_DIR / f"{msg_id[:20]}.json").write_text(
                    json.dumps({"id": msg_id, "subject": subject, "from": sender,
                                "date": date_str, "body_preview": body[:500]}, indent=2),
                    encoding="utf-8",
                )

                store_email(msg_id, subject, sender, date_str, body, collection)

                attachments = get_attachments(token, msg_id)
                if attachments is None:
                    # Fetch failed -- do NOT mark seen, so this message is
                    # retried next run instead of its attachments (which
                    # could be a survey, DD report, etc.) being silently
                    # skipped forever. The body above was still stored
                    # successfully, so store_email's upsert will just
                    # overwrite the same content again next run -- no harm.
                    log.error(f"  Could not fetch attachments for '{subject[:50]}' — will retry next run")
                    errors += 1
                    continue

                for att in attachments:
                    try:
                        route_attachment(token, msg_id, att, subject, collection)
                    except Exception as e:
                        log.error(f"  Attachment error ({att['name']}): {e}")

                seen_ids.add(msg_id)
                new += 1

            except Exception as e:
                log.error(f"Failed processing '{subject[:50]}': {e}")
                errors += 1

    save_registry(seen_ids)
    log.info(f"Email run complete — {new} new, {skipped} skipped, {errors} errors")
